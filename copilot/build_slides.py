"""Generate DEMO_SLIDES.pptx — Early Submission demo deck (v2).

Focus per user:
  - First slides: what shipped between MVP (Tue 04-28) and Early (Thu 04-30)
  - Detailed treatment of the eval framework and Langfuse trace pipeline
  - Closing slide: future-work plan to make pseudonymization & verification
    visible in the UI (the plan written tonight)

Run:  python3 build_slides.py
Out:  /Users/rikki/Desktop/Doc/OOD/openemr/copilot/DEMO_SLIDES.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).parent / "DEMO_SLIDES.pptx"

# Palette
NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
ACCENT = RGBColor(0xE8, 0x6A, 0x33)
GRAY_BG = RGBColor(0xF4, 0xF6, 0xF8)
GRAY_LINE = RGBColor(0xC9, 0xD2, 0xDA)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x60, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
PURPLE = RGBColor(0x6E, 0x40, 0xC9)


# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers --------------------------------------------------------

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_notes(slide, text: str) -> None:
    """Attach speaker notes (teleprompter copy) to a slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(12)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def slide_header(slide, eyebrow, title, page=None):
    add_rect(slide, 0, 0, SW, Inches(0.18), TEAL)
    add_text(slide, Inches(0.5), Inches(0.3), Inches(10),
             Inches(0.35), eyebrow.upper(),
             size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.5), Inches(0.6), Inches(12.3),
             Inches(0.9), title,
             size=28, bold=True, color=NAVY)
    add_rect(slide, Inches(0.5), Inches(1.45), Inches(1.0),
             Inches(0.05), ACCENT)
    if page is not None:
        add_text(slide, Inches(12.4), Inches(7.05), Inches(0.8),
                 Inches(0.3), str(page),
                 size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8),
             Inches(0.3), "Clinical Co-Pilot · Early Submission Demo",
             size=10, color=TEXT_MUTED)


def bullet_list(slide, x, y, w, h, items, *, size=14, leading=0.42,
                indent_marker="▸"):
    line_h = Inches(leading)
    for i, item in enumerate(items):
        row_y = y + i * line_h
        add_text(slide, x, row_y, Inches(0.3), line_h, indent_marker,
                 size=size, bold=True, color=ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.32), row_y, w - Inches(0.32),
                 line_h, item, size=size, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)


def two_col(slide, left_title, left_items, right_title, right_items,
            top=Inches(1.85), height=Inches(5.0)):
    col_w = Inches(6.0)
    gap = Inches(0.33)
    left_x = Inches(0.5)
    right_x = left_x + col_w + gap

    for x, title, items in [
        (left_x, left_title, left_items),
        (right_x, right_title, right_items),
    ]:
        add_rect(slide, x, top, col_w, height, GRAY_BG)
        add_rect(slide, x, top, Inches(0.18), height, TEAL)
        add_text(slide, x + Inches(0.35), top + Inches(0.18),
                 col_w - Inches(0.5), Inches(0.4), title,
                 size=15, bold=True, color=NAVY)
        bullet_list(slide, x + Inches(0.35), top + Inches(0.7),
                    col_w - Inches(0.5), height - Inches(1.0), items,
                    size=12, leading=0.42)


def code_block(slide, x, y, w, h, lines, size=12):
    add_rect(slide, x, y, w, h, RGBColor(0x14, 0x1E, 0x2E))
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                   w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = "Consolas"
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xDC, 0xE8, 0xF4)


def table(slide, x, y, col_widths, headers, rows, *,
          row_h=0.42, header_h=0.42, font_size=11):
    cur_x = x
    for i, hdr in enumerate(headers):
        add_rect(slide, cur_x, y, col_widths[i], Inches(header_h), NAVY)
        add_text(slide, cur_x + Inches(0.12), y, col_widths[i] - Inches(0.2),
                 Inches(header_h), hdr,
                 size=font_size + 1, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        cur_x += col_widths[i]
    for r, row in enumerate(rows):
        cur_x = x
        bg = WHITE if r % 2 == 0 else GRAY_BG
        for i, cell in enumerate(row):
            add_rect(slide, cur_x, y + Inches(header_h) + r * Inches(row_h),
                     col_widths[i], Inches(row_h), bg, line=GRAY_LINE)
            add_text(slide, cur_x + Inches(0.12),
                     y + Inches(header_h) + r * Inches(row_h),
                     col_widths[i] - Inches(0.2), Inches(row_h), cell,
                     size=font_size, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
            cur_x += col_widths[i]


# ---------- slides ---------------------------------------------------------

def slide_title():
    s = add_slide()
    set_notes(s, (
        "Demo deck for the Early Submission, Thursday April 30. "
        "This walks through what shipped between the MVP on Tuesday and tonight, "
        "deep-dives the eval framework and Langfuse trace pipeline, "
        "and closes on the visibility plan we'll execute on Sunday. "
        "Twenty slides total. Aim for under 8 minutes of voiceover."
    ))
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.4), SW, Inches(0.06), ACCENT)
    add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(0.5),
             "EARLY SUBMISSION · 2026-04-30",
             size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(1.2),
             "Clinical Co-Pilot",
             size=48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.7),
             "From MVP to Early Submission · Eval framework + Langfuse",
             size=20, color=RGBColor(0xC9, 0xD2, 0xDA))
    add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
             "ARCHITECTURE.md · USERS.md · AUDIT.md → /copilot/ agent service · embedded into OpenEMR",
             size=12, color=RGBColor(0x9A, 0xA8, 0xB5))


def slide_recap_mvp():
    s = add_slide()
    set_notes(s, (
        "Tuesday's MVP was three written documents totaling roughly eleven hundred lines. "
        "AUDIT.md called out the structural risks in the OpenEMR codebase — call-site-only "
        "ACL, raw PHI flowing into LLM prompts, audit-log gaps. "
        "USERS.md scoped the product to one persona, a primary care physician, with three "
        "named use cases — pre-visit brief, multi-condition reasoning, medication safety. "
        "ARCHITECTURE.md was the design we'd build against. "
        "The important point on this slide: nothing was running yet. Everything was on paper."
    ))
    slide_header(s, "Where we started", "MVP — Tuesday 04-28: docs only", 2)
    add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.4),
             "MVP shipped three docs, no running code. Architecture was on paper.",
             size=15, color=TEXT_MUTED)
    cards = [
        ("AUDIT.md", "404 lines",
         "Codebase audit. Identified call-site-only ACL, raw PHI flow risk, audit gaps."),
        ("USERS.md", "153 lines",
         "Target user (PCP). 3 use cases (UC1 brief, UC2 reasoning, UC3 med safety)."),
        ("ARCHITECTURE.md", "592 lines",
         "Full design — agent, tools, PHI minimizer, verification, Langfuse, evals."),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    top = Inches(2.4)
    for i, (title, sub, body) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(3.0), GRAY_BG)
        add_rect(s, x, top, card_w, Inches(0.5), NAVY)
        add_text(s, x + Inches(0.25), top, card_w - Inches(0.5),
                 Inches(0.5), title,
                 size=15, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.25), top + Inches(0.65),
                 card_w - Inches(0.5), Inches(0.4), sub,
                 size=12, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), top + Inches(1.05),
                 card_w - Inches(0.5), Inches(1.8), body,
                 size=12, color=TEXT_DARK)
    add_text(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
             "→ Architecture was on paper. No code, no service, no traces, no tests, no UI.",
             size=14, bold=True, color=ACCENT)


def slide_delta_overview():
    s = add_slide()
    set_notes(s, (
        "Five phases shipped between Tuesday and tonight. "
        "Phases A through D map directly to architecture sections — skeleton plus OAuth, "
        "tools plus PHI minimizer, agent loop plus verification, observability plus eval. "
        "Two new things today: Phase E3 is the iframe rail — the agent is now embedded "
        "inside OpenEMR, not just a standalone sidecar. Phase E4 is real Synthea data on "
        "Railway plus all three use cases verified end to end. "
        "Today's deep dive is Phase D — eval and Langfuse — and the closing visibility plan."
    ))
    slide_header(s, "What changed", "Early Submission — what we shipped this week", 3)
    phases = [
        ("Phase A", "Skeleton · OAuth2 · FHIR roundtrip", "✅"),
        ("Phase B", "8 FHIR tools · PHI minimizer · ACL mirror", "✅"),
        ("Phase C", "Agent loop · Layer-1 + Layer-2 verification", "✅"),
        ("Phase D", "Langfuse observability · Eval suite (17/17)", "✅"),
        ("Phase E1-2", "Chat UI · Railway deploy", "✅"),
        ("Phase E3", "Iframe rail in OpenEMR (demographics.php)", "✅ NEW"),
        ("Phase E4", "10 Synthea patients · UC1/UC2/UC3 smoke tests pass", "✅ NEW"),
    ]
    top = Inches(1.85)
    chip_h = Inches(0.65)
    chip_gap = Inches(0.11)
    for i, (label, body, mark) in enumerate(phases):
        y = top + i * (chip_h + chip_gap)
        add_rect(s, Inches(0.5), y, Inches(12.3), chip_h, GRAY_BG)
        color = ACCENT if "NEW" in mark else TEAL
        add_rect(s, Inches(0.5), y, Inches(0.18), chip_h, color)
        add_text(s, Inches(0.85), y, Inches(2.0), chip_h, label,
                 size=14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(2.85), y, Inches(8.0), chip_h, body,
                 size=13, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(11.0), y, Inches(1.7), chip_h, mark,
                 size=14, bold=True, color=GREEN if "NEW" not in mark else ACCENT,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
             "Today's deep dive: Phase D (eval + Langfuse) and the Phase F visibility plan.",
             size=11, color=TEXT_MUTED)


def slide_architecture_traceback():
    s = add_slide()
    set_notes(s, (
        "Every section in ARCHITECTURE.md has a corresponding directory under copilot/app/. "
        "If you wrote it, you can find it. "
        "Two notes worth calling out: Layer-2 says two-of-four — we ship cross-patient "
        "leakage and allergy contraindication this week, with renal-dose and QTc queued "
        "for Sunday Final. The bottom row is the iframe rail integration added tonight."
    ))
    slide_header(s, "From design to code",
                 "Every architecture section has a corresponding directory", 4)
    headers = ["ARCHITECTURE.md §", "Component", "Where it lives", "Status"]
    rows = [
        ("§2.1 Custom orchestration", "Tool-calling loop (OpenAI + Anthropic)", "app/agent/loop.py", "✅"),
        ("§3.1 Tools",                 "8 FHIR-backed tools",                    "app/tools/",       "✅"),
        ("§3.2 5-step pattern",        "Shared helper",                          "app/tools/_base.py:run_tool", "✅"),
        ("§3.3 PHI minimizer",         "Strip + session pseudonyms",             "app/phi/",         "✅"),
        ("§3.4 Trust boundaries",      "OAuth2 + ACL mirror",                    "app/fhir/, acl/",  "✅"),
        ("§4.1 Layer-1 verifier",      "Source attribution",                     "verification/attribution.py", "✅"),
        ("§4.1 Layer-2 verifier",      "Domain rules (allergy, leakage)",        "verification/rules.py",       "✅ (2 of 4)"),
        ("§5  Observability",          "Langfuse wrapper",                       "app/observability/trace.py",  "✅"),
        ("§6  Evaluation",             "pytest harness + RESULTS.md",            "evals/",                       "✅"),
        ("§10 Deployment",             "Railway service `copilot`",              "Dockerfile · railway.toml",    "✅"),
        ("§10 Embedding",              "iframe rail in OpenEMR",                 "demographics.php (+54 lines)", "✅ NEW"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(2.7), Inches(3.4), Inches(4.4), Inches(1.8)],
          headers, rows, row_h=0.4)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
             "One env var swaps Anthropic ↔ OpenAI. Both adapters tested. "
             "Sunday → flip back to Sonnet 4.6.",
             size=12, color=TEXT_MUTED)


def slide_iframe_rail():
    s = add_slide()
    set_notes(s, (
        "Fifty-four lines added to demographics.php — the patient summary page in OpenEMR. "
        "It looks up the patient's FHIR UUID from patient_data.uuid, builds the iframe src, "
        "and renders a 36-pixel collapsed tab on the right edge of every chart. "
        "Click that tab and a 400-pixel rail slides in. "
        "The chat UI auto-binds to the patient_id query param, so the physician never types "
        "a UUID by hand. This satisfies the PRD's requirement that the agent be embedded "
        "directly into OpenEMR. The standalone URL still works for demo and dev."
    ))
    slide_header(s, "Embedded into OpenEMR",
                 "Iframe rail integration — demographics.php → Co-Pilot tab", 5)
    # Left — diagram-ish
    add_rect(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(5.0), GRAY_BG)
    add_text(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(0.4),
             "How it works", size=15, bold=True, color=NAVY)
    code_block(s, Inches(0.7), Inches(2.5), Inches(5.6), Inches(4.2), [
        "interface/patient_file/summary/",
        "  demographics.php   (+54 lines, additive)",
        "    │",
        "    ├─ lookup patient_data.uuid → FHIR uuid",
        "    ├─ build src=https://copilot…/?patient_id=<uuid>",
        "    └─ render 36px collapsed tab on right edge",
        "",
        "Click `Co-Pilot ▸` → body.copilot-open",
        "    → 400px iframe slides in",
        "    → web/index.html reads ?patient_id=<uuid>",
        "    → auto-fires startSession() after 50ms",
        "",
        "Physician never types the FHIR id by hand.",
    ], size=11)
    # Right — what it solves
    add_text(s, Inches(6.85), Inches(1.85), Inches(6.0), Inches(0.4),
             "Why it matters", size=15, bold=True, color=NAVY)
    items = [
        "Satisfies PRD §2: \"AI agent embedded directly into OpenEMR\"",
        "Pre-scoped to active patient — zero context switching",
        "Standalone /'/' URL still served — for demo + dev",
        "Tested on local OpenEMR + deployed on Railway",
        "Pushed to GitHub master (commit 937f42a83)",
        "Auto-deploys to https://openemr-production-...up.railway.app",
    ]
    bullet_list(s, Inches(6.85), Inches(2.5), Inches(6.0),
                Inches(4.0), items, size=13, leading=0.55)


def slide_synthea_smoke():
    s = add_slide()
    set_notes(s, (
        "Ten alive Synthea patients imported on Railway OpenEMR. "
        "The previous blocker — the documents directory not being writable inside the "
        "Railway container — was fixed by flipping the write bit on sites/default. "
        "Mariela is the demo star: 47-year-old female, 33 encounters, LDL 190, on "
        "chlorpheniramine, creatinine 2.72 — perfect for UC1 and UC2. "
        "Dana is the UC3 hard-block demo: a 2-year-old with ten allergies including aspirin. "
        "The smoke table at the bottom is the actual measured result against the deployed "
        "agent — every use case has a passing trace, with verification and pseudonyms intact."
    ))
    slide_header(s, "Real data + smoke tests",
                 "10 Synthea patients on Railway · all 3 use cases verified", 6)
    # Patient roster (compact)
    headers = ["Patient", "Age/Sex", "Encounters", "Best for"]
    rows = [
        ("Mariela",   "47F", "33", "UC1 + UC2 — LDL 190, chlorpheniramine, creatinine 2.72"),
        ("Dana",      "2y",  "17", "UC3 hard-block — 10 allergies including aspirin"),
        ("Kacie",     "30F", "39", "Clean baseline"),
        ("Kiera",     "—",   "30", "Clean baseline"),
        ("Un",        "62F", "39", "Clean baseline"),
        ("Qiana",     "39F", "50", "Clean baseline"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(1.6), Inches(1.2), Inches(1.6), Inches(7.0)],
          headers, rows, row_h=0.36, font_size=11)

    # Smoke results
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
             "Use-case smoke tests against deployed agent",
             size=14, bold=True, color=NAVY)
    sheaders = ["UC", "Patient", "Result"]
    srows = [
        ("UC1", "Mariela", "5 cited claims · verification_passed=true · 12-15s · pseudonym Patient-Z3DU"),
        ("UC2", "Mariela", "2 claims · renal-aware reasoning · honestly flags `no recent vitals` · 8s"),
        ("UC3", "Dana + aspirin", "Layer-2 hard-block fires cleanly · refusal returned · 5-7s"),
    ]
    table(s, Inches(0.5), Inches(5.45),
          [Inches(0.9), Inches(2.0), Inches(9.5)],
          sheaders, srows, row_h=0.42, font_size=11)


# ============================================================
#  EVAL — DETAILED
# ============================================================

def slide_eval_intro():
    s = add_slide()
    set_notes(s, (
        "Eval is the offline truth. Seventeen tests, all passing. "
        "Four design principles to call out as you read the right column. "
        "First: same gate code as production — we don't stub the verifier in tests. "
        "Second: hand-crafted FHIR fixtures, not Synthea, because eval needs determinism. "
        "Third: Synthea is for the demo, fixtures are for eval — different jobs. "
        "Fourth: RESULTS.md is regenerated by pytest_sessionfinish on every run, "
        "so the repo's eval status is never out of date."
    ))
    slide_header(s, "Eval framework",
                 "How we know the agent is right — before it ships", 7)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.6),
             "Eval is the offline truth. Runs in CI on every change, "
             "asserts properties of TurnTrace, blocks releases on regression.",
             size=15, color=TEXT_DARK)
    # Big number
    add_rect(s, Inches(0.5), Inches(2.9), Inches(4.0), Inches(3.6), NAVY)
    add_text(s, Inches(0.5), Inches(3.0), Inches(4.0), Inches(0.5),
             "TESTS PASSING",
             size=12, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.5), Inches(4.0), Inches(2.0),
             "17 / 17",
             size=72, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.4), Inches(4.0), Inches(0.4),
             "0 failed · 0 skipped",
             size=14, color=RGBColor(0xC9, 0xD2, 0xDA),
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.8), Inches(4.0), Inches(0.5),
             "(`make eval` exit 0)",
             size=11, color=RGBColor(0x9A, 0xA8, 0xB5),
             align=PP_ALIGN.CENTER)
    # Right column principles
    items = [
        "Same gate code as production — no test-doubles for the verifier",
        "Hand-crafted FHIR fixtures, deterministic & fast",
        "Synthea = demo data; fixtures = eval data (different jobs)",
        "RESULTS.md auto-regenerated each run — repo status never lies",
        "live_llm marker: real-LLM tests behind ANTHROPIC_LIVE=1",
    ]
    add_text(s, Inches(4.85), Inches(2.9), Inches(8.0), Inches(0.4),
             "Design principles", size=15, bold=True, color=NAVY)
    bullet_list(s, Inches(4.85), Inches(3.35), Inches(8.0),
                Inches(3.2), items, size=13, leading=0.55)


def slide_eval_categories():
    s = add_slide()
    set_notes(s, (
        "Three categories. "
        "Factual attribution catches hallucinations — every clinical claim must cite a "
        "record_id from a real tool result. Unanchored claims get stripped. "
        "Adversarial — cross-patient leakage, prompt injection, ACL bypass — is the "
        "non-negotiable 100% pass gate. Anything below 100% in this row blocks the release. "
        "Failure-mode handling makes sure the agent refuses on missing data instead of "
        "guessing. The 17/17 number is meaningful only because the adversarial column is "
        "also all green."
    ))
    slide_header(s, "Eval framework",
                 "Three categories · two are non-negotiable hard gates", 8)
    headers = ["Category", "What it checks", "Example tests (file::test)"]
    rows = [
        ("Factual / attribution",
         "Every clinical claim cites a record_id from a real tool result. Strips claims with no tool source.",
         "test_verification.py::test_attribution_passes_when_all_claims_anchored\n"
         "test_verification.py::test_attribution_strips_unanchored_claim"),
        ("Adversarial  ★ 100% pass",
         "Cross-patient leakage · prompt injection · ACL bypass — anything below 100% blocks release.",
         "test_verification.py::test_cross_patient_leakage_hard_blocks\n"
         "test_scenarios.py::test_prompt_injection_does_not_leak_other_patients\n"
         "test_tool_integration.py::test_acl_denies_unknown_role"),
        ("Failure-mode handling",
         "Refuses on missing data · blocks contraindicated meds · honest data_gaps.",
         "test_scenarios.py::test_uc1_refuses_when_no_prior_encounter\n"
         "test_verification.py::test_allergy_contraindication_blocks_safe_verdict"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(3.0), Inches(4.5), Inches(4.8)],
          headers, rows, row_h=1.2, font_size=11)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
             "★ Cross-patient leakage and ACL bypass are hard 100%-pass gates "
             "in evals/RESULTS.md — the 17/17 number isn't an aggregate, the adversarial column must be all-green.",
             size=12, color=ACCENT, bold=True)


def slide_eval_test_inventory():
    s = add_slide()
    set_notes(s, (
        "This is what the seventeen tests actually do, broken out by file. "
        "Seven PHI minimizer tests prove that identifiers come out — name, telecom, "
        "address, identifier, birthDate — and clinical codes stay in — RxNorm, LOINC, "
        "ICD-10. They also prove pseudonyms are stable within a session. "
        "Three tool integration tests: record_id capture, ACL deny path, RxNorm capture. "
        "Four verification tests are the spine of the safety story — Layer-1 anchored "
        "vs unanchored, Layer-2 cross-patient leakage hard-block, Layer-2 allergy "
        "contraindication block. "
        "Three live-LLM scenarios run only with ANTHROPIC_LIVE=1: UC1 happy path, "
        "UC1 refuses on no-prior-encounter, prompt injection cannot leak another patient."
    ))
    slide_header(s, "Eval framework",
                 "Full inventory — what each of the 17 tests actually does", 9)
    headers = ["File", "# tests", "What it asserts"]
    rows = [
        ("evals/tools/test_phi_minimizer.py", "7",
         "Strip drops identifiers · keeps RxNorm/LOINC/ICD-10 · pseudonym stable within session · provider name pseudonymized"),
        ("evals/tools/test_tool_integration.py", "3",
         "Tool returns record_ids · ACL deny path returns no data · RxNorm captured for medications"),
        ("evals/agent/test_verification.py", "4",
         "L1 attribution — anchored claims pass · unanchored claims stripped · L2 cross-patient leakage hard-blocks · L2 allergy contraindication blocks 'safe' verdicts"),
        ("evals/agent/test_scenarios.py", "3",
         "Live-LLM (under ANTHROPIC_LIVE=1): UC1 happy path anchors every claim · UC1 refuses on no-prior-encounter · prompt injection cannot leak another patient"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(4.0), Inches(1.0), Inches(7.3)],
          headers, rows, row_h=0.95, font_size=11)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             "Each line maps to ≥1 category from the previous slide. "
             "The 4 verification-gate tests + 3 adversarial scenarios are the spine of the safety story.",
             size=12, color=TEXT_MUTED, bold=True)


def slide_eval_run_flow():
    s = add_slide()
    set_notes(s, (
        "Three Make targets. "
        "make test runs the offline subset — PHI plus tool tests, no LLM. "
        "make eval runs the full suite with mocked LLM. "
        "make eval-live runs with ANTHROPIC_LIVE=1 and hits the real API. "
        "Two pytest hooks in conftest.py make this work: the live_llm marker auto-skips "
        "real-LLM tests by default, and pytest_sessionfinish writes RESULTS.md after each "
        "run. The bottom strip is a verbatim sample of what RESULTS.md actually contains."
    ))
    slide_header(s, "Eval framework",
                 "Layout · how it runs · what RESULTS.md produces", 10)

    # Left — file tree
    add_rect(s, Inches(0.5), Inches(1.85), Inches(6.1), Inches(4.2), GRAY_BG)
    add_text(s, Inches(0.7), Inches(1.95), Inches(6), Inches(0.4),
             "evals/  layout",
             size=14, bold=True, color=NAVY)
    code_block(s, Inches(0.7), Inches(2.4), Inches(5.7), Inches(3.55), [
        "evals/",
        "├── conftest.py             # markers + RESULTS.md autogen",
        "├── RESULTS.md              # auto-written after every run",
        "├── tools/",
        "│   ├── test_phi_minimizer.py     7 tests",
        "│   └── test_tool_integration.py  3 tests",
        "└── agent/",
        "    ├── test_verification.py      4 tests",
        "    └── test_scenarios.py         3 live-LLM",
    ], size=11)

    # Right — commands + flow
    add_rect(s, Inches(6.85), Inches(1.85), Inches(6.0), Inches(4.2), GRAY_BG)
    add_text(s, Inches(7.05), Inches(1.95), Inches(6), Inches(0.4),
             "How it runs", size=14, bold=True, color=NAVY)
    code_block(s, Inches(7.05), Inches(2.4), Inches(5.6), Inches(1.4), [
        "$ make test         # PHI + tool tests, no LLM",
        "$ make eval         # full suite, mocked LLM",
        "$ make eval-live    # ANTHROPIC_LIVE=1, real API",
    ], size=11)
    add_text(s, Inches(7.05), Inches(3.95), Inches(5.6), Inches(0.4),
             "Hooks (conftest.py)", size=13, bold=True, color=NAVY)
    bullets = [
        "live_llm marker → skips real-LLM tests by default",
        "pytest_sessionfinish → writes evals/RESULTS.md",
        "RESULTS.md never hand-edited — repo status is honest",
    ]
    bullet_list(s, Inches(7.05), Inches(4.4), Inches(5.6), Inches(1.4),
                bullets, size=12, leading=0.5)

    # Bottom — RESULTS.md teaser
    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.95), NAVY)
    code_block(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.85), [
        "# Eval Suite Results — Generated by `make eval`. Exit status: 0",
        "# ✅ Passed: 17 · ❌ Failed: 0 · ⏭ Skipped: 0",
    ], size=11)


# ============================================================
#  LANGFUSE — DETAILED
# ============================================================

def slide_langfuse_intro():
    s = add_slide()
    set_notes(s, (
        "Langfuse is the online story — what the agent actually did at runtime. "
        "About ninety lines in one file. Two-tier design: real Langfuse if keys are "
        "present, no-op fallback otherwise — falls through to structured stdout. "
        "The agent itself never imports Langfuse — only main.py does. That's deliberate "
        "decoupling. Observability can never crash the agent: every exception in emit() "
        "is caught and warned, never raised. Same TurnTrace eval asserts on — "
        "one canonical telemetry object across both pipelines."
    ))
    slide_header(s, "Langfuse observability",
                 "How we know what the agent did — after it ran", 11)
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.6),
             "Online forensics for the agent. One file (~90 lines). "
             "Fire-and-forget. Cannot break the agent — ever.",
             size=15, color=TEXT_DARK)

    items_left = [
        "One file: app/observability/trace.py",
        "Two-tier design — real Langfuse if keys present, no-op fallback otherwise",
        "Lazy import — Langfuse SDK never loaded in noop path",
        "All exceptions in emit() are caught + warned, never raised",
    ]
    items_right = [
        "Agent never imports Langfuse — only main.py does",
        "Loop builds TurnTrace; tracer.emit() ships it",
        "Same TurnTrace eval asserts on (one canonical telemetry object)",
        "Fall-back: structured stdout — Railway-friendly",
    ]
    two_col(s, "Design rules", items_left,
            "Decoupling", items_right, top=Inches(2.6), height=Inches(4.0))


def slide_langfuse_what_is_a_turn():
    s = add_slide()
    set_notes(s, (
        "A turn is one POST to /v1/chat — physician question in, verified response out. "
        "Inside that turn the agent might call three or four tools, run the verifier, "
        "and assemble a response — but every one of those internal steps is a span inside "
        "one trace, not a separate trace. "
        "One physician question → one Langfuse record. That's the unit of observation."
    ))
    slide_header(s, "Langfuse observability",
                 "What's a 'turn'? — one POST /v1/chat round trip", 12)

    # Diagram
    add_rect(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.5), GRAY_BG)
    code_block(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(4.2), [
        "Browser  POST /v1/chat                                    ┐",
        "    {  session_id: 's_42',                                │",
        "       question:   'Brief me on this patient' }           │  ONE TURN STARTS",
        "                                                          │",
        "  ┌── agent loop iterates ─────────────────────────────────┤",
        "  │   • get_patient_summary                                │",
        "  │   • get_active_medications                             │  inner steps =",
        "  │   • get_recent_labs                                    │  Langfuse SPANS",
        "  │   • submit_response                                    │  (NOT separate turns)",
        "  └────────────────────────────────────────────────────────┤",
        "                                                          │",
        "  Layer-1 attribution + Layer-2 rules                     │",
        "                                                          │",
        "Response  { prose, claims[], data_gaps[],   trace }       │  ONE TURN ENDS",
        "                                                          │",
        "tracer.emit(trace, response)        ─────────────────── ─ ┘  ONE LANGFUSE TRACE",
    ], size=11)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "One physician question → one Langfuse trace, with all internal tool calls + token usage + verification verdict packed into its metadata.",
             size=13, color=TEXT_MUTED, bold=True)


def slide_langfuse_request_response():
    s = add_slide()
    set_notes(s, (
        "What actually ships per turn. "
        "Request side, top left: input is the natural-language question plus pseudonymous "
        "identity — user_id, session_id, patient pseudonym. Never a real patient ID. "
        "Response side, top right: output is the verified AgentResponse — prose plus "
        "claims with their record_ids plus data gaps. "
        "Bottom block is metadata and usage. Metadata holds the tool sequence, per-tool "
        "latencies, and verification verdict. Usage holds the token counts including "
        "cache_read_input — that's what proves prompt caching is hitting."
    ))
    slide_header(s, "Langfuse observability",
                 "What ships per turn — request side and response side", 13)
    # Left: input (request side)
    add_text(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(0.4),
             "Request side  (input + identity)",
             size=14, bold=True, color=NAVY)
    code_block(s, Inches(0.5), Inches(2.3), Inches(6.0), Inches(2.2), [
        "input = { 'question': trace.question_text }",
        "",
        "user_id           = trace.user_id",
        "session_id        = trace.session_id",
        "patient_pseudonym = 'Patient-a7b3'   ← never real ID",
    ], size=12)
    # Right: output (response side)
    add_text(s, Inches(6.85), Inches(1.85), Inches(6.0), Inches(0.4),
             "Response side  (verified output)",
             size=14, bold=True, color=NAVY)
    code_block(s, Inches(6.85), Inches(2.3), Inches(6.0), Inches(2.2), [
        "output = response   # post-verification dict:",
        "",
        "{ prose:     '...',",
        "  claims:    [{text, record_id}, ...],",
        "  data_gaps: ['No BP recorded in last 90 days'] }",
    ], size=12)
    # Bottom: metadata + usage
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
             "Metadata + usage  (everything else from the turn)",
             size=14, bold=True, color=NAVY)
    code_block(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(2.0), [
        "metadata = {",
        "  'tool_call_sequence':      ['get_patient_summary', 'get_active_medications', ...],",
        "  'tool_latencies_ms':       { 'get_patient_summary': 412, ... },",
        "  'tool_failures':           {},",
        "  'verification_passed':     True,",
        "  'verification_rejections': [...],   # Layer-1 unknown record_ids",
        "  'domain_rule_rejections':  [...],   # Layer-2 reasons",
        "}",
        "usage = { 'input': 3214, 'output': 487, 'cache_read_input': 2891 }",
    ], size=11)


def slide_langfuse_touchpoints():
    s = add_slide()
    set_notes(s, (
        "Three lines in main.py — that's the entire Langfuse surface area in the codebase. "
        "Line 15: lazy import. Line 23: build the tracer once at FastAPI startup. "
        "Line 107: emit once per chat request. "
        "The bottom block is the get_tracer fallback — if both Langfuse keys are present "
        "and the SDK initializes, you get a real tracer; otherwise you get a no-op that "
        "logs to stdout. This makes the wiring work in CI and on Railway without keys "
        "and lets us run locally without a Langfuse instance."
    ))
    slide_header(s, "Langfuse observability",
                 "Exactly three touch points — that's the entire surface area", 14)

    code_block(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(2.6), [
        "# main.py:15      lazy import",
        "from app.observability.trace import get_tracer",
        "",
        "# main.py:23      one tracer per process, built at FastAPI startup",
        "app.state.tracer = get_tracer(settings)",
        "",
        "# main.py:107     called once per /v1/chat request",
        "app.state.tracer.emit(output.trace, payload['response'])",
    ], size=14)

    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
             "get_tracer() — two-tier fallback",
             size=14, bold=True, color=NAVY)
    code_block(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), [
        "def get_tracer(settings):",
        "    if settings.langfuse_public_key and settings.langfuse_secret_key:",
        "        try:    return LangfuseTracer(settings)",
        "        except: pass",
        "    return _NoopTracer()    # stdout structured logs",
    ], size=12)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             "Decouples 'we always emit a trace' from 'Langfuse must be reachable'.",
             size=12, color=TEXT_MUTED, bold=True)


def slide_langfuse_phi_split():
    s = add_slide()
    set_notes(s, (
        "The orange callout at the top is the file's docstring quoted verbatim. "
        "Pseudonyms only in Langfuse. Real IDs go to OpenEMR's audit table — the "
        "HIPAA-mandated trail. "
        "Two logs, two audiences. Engineers read Langfuse to debug latency or cost. "
        "Compliance reads the audit table when an audit happens. They have different "
        "retention requirements and different threat models. Trace data and HIPAA audit "
        "data never co-mingle."
    ))
    slide_header(s, "Langfuse observability",
                 "The PHI rule — pseudonyms only, never real IDs", 15)

    # The rule (quoted)
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.4), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(1.85), Inches(0.18), Inches(1.4), ACCENT)
    add_text(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(1.2),
             "\"We never log raw PHI. Pseudonyms only. The clinical audit log "
             "(separate write to OpenEMR's existing audit table) carries the "
             "user/patient/time triple for HIPAA — Langfuse holds the technical trace.\"\n"
             "    — copilot/app/observability/trace.py docstring",
             size=13, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Two-log split table
    headers = ["", "Langfuse", "OpenEMR audit table"]
    rows = [
        ("Audience",       "Engineers",                                    "Compliance"),
        ("Contents",       "Tool traces, latency, tokens, verdicts, pseudonyms", "user_id, patient_id, timestamp, action — real IDs"),
        ("Retention",      "Engineering choice",                           "HIPAA-mandated"),
        ("Hosted",         "Self-hosted Langfuse (planned)",               "OpenEMR's existing infra"),
    ]
    table(s, Inches(0.5), Inches(3.5),
          [Inches(2.0), Inches(4.6), Inches(5.7)],
          headers, rows, row_h=0.5, font_size=11)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
             "Two logs. Two audiences. Trace data and HIPAA audit data never co-mingle.",
             size=14, bold=True, color=ACCENT)


# ============================================================
#  WORKFLOW + EVAL/LANGFUSE INTEGRATION
# ============================================================

def slide_workflow():
    s = add_slide()
    set_notes(s, (
        "End-to-end through one turn. "
        "Question comes in via /v1/chat. The agent loop builds a TurnTrace as it runs — "
        "every tool call appends to the sequence, every token count gets added in, "
        "every verifier rejection gets recorded. "
        "Then main.py emits the trace. Two destinations in parallel: Langfuse for engineers "
        "with pseudonyms, OpenEMR's audit table for compliance with real IDs. "
        "Finally the verified response goes back to the physician with cited record_ids."
    ))
    slide_header(s, "End-to-end workflow",
                 "One chat turn — eval and Langfuse share the same TurnTrace", 16)
    box_w = Inches(11.0)
    box_x = Inches(1.15)
    y = Inches(1.85)
    steps = [
        ("Physician asks question via /v1/chat", NAVY),
        ("Agent loop builds TurnTrace as it runs:\n"
         "ACL check → tool sequence → token counts → Layer-1 + Layer-2 verification",
         TEAL),
        ("Verifier sets verification_passed; "
         "rejections → verification_rejections / domain_rule_rejections",
         TEAL),
        ("main.py:107  →  tracer.emit(trace, response)", ACCENT),
        ("→ Langfuse (engineering trace, pseudonyms)\n"
         "→ OpenEMR audit table (compliance, real IDs · separate write)",
         NAVY),
        ("Verified response with cited record_ids returned to physician", GREEN),
    ]
    h = Inches(0.78)
    gap = Inches(0.1)
    for i, (txt, color) in enumerate(steps):
        yy = y + i * (h + gap)
        add_rect(s, box_x, yy, box_w, h, GRAY_BG)
        add_rect(s, box_x, yy, Inches(0.18), h, color)
        add_text(s, box_x + Inches(0.35), yy, box_w - Inches(0.5), h, txt,
                 size=13, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)


def slide_eval_vs_langfuse():
    s = add_slide()
    set_notes(s, (
        "Same TurnTrace object. Eval asserts properties of it offline, Langfuse ships it "
        "online. There's no parallel logging code path that could drift between test and "
        "production. "
        "Eval blocks the release on regression. Langfuse never blocks anything — "
        "observability is fire-and-forget. "
        "Two-line takeaway: eval proves the agent is right; Langfuse proves what the "
        "agent did."
    ))
    slide_header(s, "Eval ↔ Langfuse",
                 "Same TurnTrace · two lifecycles · no parallel logging path", 17)
    headers = ["", "Eval (offline)", "Langfuse (online)"]
    rows = [
        ("When",        "CI on every change",          "Production, every chat turn"),
        ("Input",       "Hand-crafted FHIR fixtures",  "Synthea patient traffic"),
        ("Asserts on",  "TurnTrace properties",        "Ships TurnTrace as a span"),
        ("Catches",     "Regressions before deploy",   "Drift / latency / cost after deploy"),
        ("Fails how",   "Blocks the release",          "Never blocks the agent"),
        ("Output",      "RESULTS.md + exit code",      "Langfuse UI + structured stdout"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(2.5), Inches(4.7), Inches(5.1)],
          headers, rows, row_h=0.5, font_size=12)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.2),
             "One canonical telemetry object · two lifecycles · no parallel logging path to drift.\n\n"
             "Eval proves the agent is right · Langfuse proves what the agent did.",
             size=15, bold=True, color=ACCENT)


# ============================================================
#  FUTURE WORK — VISIBILITY PLAN
# ============================================================

def slide_future_visibility_overview():
    s = add_slide()
    set_notes(s, (
        "Reviewer feedback this week: the two load-bearing security features — PHI "
        "pseudonymization and the verification gate — happen invisibly. The viewer of "
        "the UI sees only the final prose plus a 'verified' pill. They can't see "
        "what was stripped or which claims passed which layer. "
        "The Sunday plan is three new collapsible panels per response. "
        "PHI panel: counts plus field paths plus pseudonyms. Never any actual values. "
        "Verification timeline: Layer-1 plus Layer-2 status, expanded when failing. "
        "Per-claim icons: accepted versus stripped, with reason on hover. "
        "Nothing leaks PHI to the browser — the rule is field paths and counts only."
    ))
    slide_header(s, "Future work · Sunday Final",
                 "Make pseudonymization & verification visible in the UI", 18)

    # Header callout
    add_rect(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.0), NAVY)
    add_text(s, Inches(0.7), Inches(1.95), Inches(12.0), Inches(0.8),
             "Reviewer feedback: the two load-bearing security features happen "
             "invisibly. Plan: make both observable per response, without ever "
             "leaking PHI to the browser.",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, bold=True)

    # Three-pillar plan
    pillars = [
        ("🔒 PHI Minimization panel",
         "Per-response collapsible · counts + field paths + pseudonyms.",
         "Patient (Patient-A1B2): name(14), telecom, birthDate→age, identifier\n"
         "MedicationRequest (×3): requester→Provider-C, RxNorm preserved\n"
         "Free-text scrub: 2 name replacements"),
        ("✅ Verification timeline",
         "Layer-1 + Layer-2 status, expanded when failing.",
         "Layer 1 (Attribution): 5/5 claims anchored\n"
         "Layer 2 (Domain rules): no contraindications\n"
         "or → 🛑 Layer 2: BLOCKED — allergy_contraindication"),
        ("⚠️ Per-claim status icons",
         "✅ accepted · ⚠️ stripped (reason on hover).",
         "Adds a 'Stripped claims' subsection that only renders when non-empty.\n"
         "Demo trick: prompt-injection lights up the strip-and-explain path."),
    ]
    top = Inches(3.05)
    card_w = Inches(4.0)
    gap = Inches(0.15)
    start_x = Inches(0.5)
    for i, (title, sub, body) in enumerate(pillars):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, top, card_w, Inches(3.7), GRAY_BG)
        add_rect(s, x, top, card_w, Inches(0.55), PURPLE)
        add_text(s, x + Inches(0.2), top, card_w - Inches(0.4),
                 Inches(0.55), title,
                 size=13, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), top + Inches(0.65),
                 card_w - Inches(0.4), Inches(0.5), sub,
                 size=11, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.2), top + Inches(1.15),
                 card_w - Inches(0.4), Inches(2.5), body,
                 size=11, color=TEXT_DARK)


def slide_future_visibility_plan():
    s = add_slide()
    set_notes(s, (
        "Ten ordered steps. The key design move is step two: the sibling-inspector pattern. "
        "We do not change the strip_* function signatures. Instead we add a parallel "
        "report_* module that walks the same raw FHIR input and emits a structured report "
        "of what would be removed — without ever inspecting the stripped output. "
        "This means zero churn in the existing seventeen PHI tests. We add two new tests "
        "on top to cover the new reports — including a regression guard that asserts no "
        "raw value from the input ever appears in the report's serialization. "
        "Steps 7 and 8 are the wire-up — fold the new fields into TurnTrace, render three "
        "new panels in the existing chat UI."
    ))
    slide_header(s, "Future work · Sunday Final",
                 "Implementation outline — sibling-inspector pattern (zero test churn)", 19)

    headers = ["#", "Step", "Files"]
    rows = [
        ("1", "Add MinimizationReport, ToolCallTrace, RuleHit, ClaimVerification schemas", "app/agent/schemas.py"),
        ("2", "New module: sibling inspectors (no signature changes to strip_*)", "app/phi/report.py (NEW)"),
        ("3", "Add PseudonymMap.summary() — counts only", "app/phi/session.py"),
        ("4", "Tool base: capture inspector output + acl_allowed + duration_ms", "app/tools/_base.py"),
        ("5", "Verifier: expose accepted_claims alongside rejected", "verification/attribution.py"),
        ("6", "Rules: list[str] → list[RuleHit] (rule, record_id, message)", "verification/rules.py"),
        ("7", "Loop: fold new fields into TurnTrace", "app/agent/loop.py"),
        ("8", "Frontend: 3 new collapsible <details> + per-claim icons", "app/web/index.html"),
        ("9", "+2 tests: report shape · no-raw-value-leak invariant", "evals/tools/test_phi_minimizer.py"),
        ("10", "+2 tests: accepted/rejected breakdown · L2 rule attribution", "evals/agent/test_verification.py"),
    ]
    table(s, Inches(0.5), Inches(1.85),
          [Inches(0.5), Inches(7.4), Inches(4.4)],
          headers, rows, row_h=0.4, font_size=11)
    # Bottom callout
    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), NAVY)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
             "Pseudonyms only · field paths only · char counts only — never values. "
             "Same TurnTrace shape; just stop collapsing it before the UI sees it.",
             size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_thanks():
    s = add_slide()
    set_notes(s, (
        "Thank you. The recorded demo follows this slide — UC1 brief on Mariela, "
        "UC2 dizziness reasoning, UC3 aspirin hard-block on Dana. "
        "Both deployed URLs are listed at the bottom of the slide for the reviewer."
    ))
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(3.4), SW, Inches(0.06), ACCENT)
    add_text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(1.0),
             "Thank you",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.6),
             "Demo follows — UC1 brief · UC2 reasoning · UC3 medication safety",
             size=18, color=RGBColor(0xC9, 0xD2, 0xDA))
    add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.4),
             "https://openemr-production-0c8c.up.railway.app/",
             size=13, color=ACCENT)
    add_text(s, Inches(0.8), Inches(6.85), Inches(12), Inches(0.4),
             "https://copilot-production-b532.up.railway.app/",
             size=13, color=ACCENT)


# ---------- build ----------------------------------------------------------

def main():
    slide_title()
    slide_recap_mvp()
    slide_delta_overview()
    slide_architecture_traceback()
    slide_iframe_rail()
    slide_synthea_smoke()
    slide_eval_intro()
    slide_eval_categories()
    slide_eval_test_inventory()
    slide_eval_run_flow()
    slide_langfuse_intro()
    slide_langfuse_what_is_a_turn()
    slide_langfuse_request_response()
    slide_langfuse_touchpoints()
    slide_langfuse_phi_split()
    slide_workflow()
    slide_eval_vs_langfuse()
    slide_future_visibility_overview()
    slide_future_visibility_plan()
    slide_thanks()
    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} — {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
